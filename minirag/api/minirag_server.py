from fastapi import FastAPI, HTTPException, File, UploadFile, Form, Request, BackgroundTasks

# Backend (Python)
# Add this to store progress globally
from typing import Dict
import threading

# Global progress tracker
scan_progress: Dict = {
    "is_scanning": False,
    "current_file": "",
    "indexed_count": 0,
    "total_files": 0,
    "progress": 0,
}

# Lock for thread-safe operations
progress_lock = threading.Lock()

import json
import os

from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import logging
import argparse
import time
import re
from typing import List, Dict, Any, Optional, Union
from minirag import MiniRAG, QueryParam
from minirag.api import __api_version__

from minirag.utils import EmbeddingFunc
from enum import Enum
from pathlib import Path
import shutil
import aiofiles
from ascii_colors import trace_exception, ASCIIColors
import sys
import configparser

from fastapi import Depends, Security
from fastapi.security import APIKeyHeader
from fastapi.middleware.cors import CORSMiddleware
from contextlib import asynccontextmanager

from starlette.status import HTTP_403_FORBIDDEN
import pipmaster as pm

from dotenv import load_dotenv

# 加载环境变量
load_dotenv("config.env")


def estimate_tokens(text: str) -> int:
    """Estimate the number of tokens in text
    Chinese characters: approximately 1.5 tokens per character
    English characters: approximately 0.25 tokens per character
    """
    # Use regex to match Chinese and non-Chinese characters separately
    chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
    non_chinese_chars = len(re.findall(r"[^\u4e00-\u9fff]", text))

    # Calculate estimated token count
    tokens = chinese_chars * 1.5 + non_chinese_chars * 0.25

    return int(tokens)


class OllamaServerInfos:
    # Constants for emulated Ollama model information
    LIGHTRAG_NAME = "minirag"
    LIGHTRAG_TAG = os.getenv("OLLAMA_EMULATING_MODEL_TAG", "latest")
    LIGHTRAG_MODEL = f"{LIGHTRAG_NAME}:{LIGHTRAG_TAG}"
    LIGHTRAG_SIZE = 7365960935  # it's a dummy value
    LIGHTRAG_CREATED_AT = "2024-01-15T00:00:00Z"
    LIGHTRAG_DIGEST = "sha256:minirag"

    KV_STORAGE = "JsonKVStorage"
    DOC_STATUS_STORAGE = "JsonDocStatusStorage"
    GRAPH_STORAGE = "NetworkXStorage"
    VECTOR_STORAGE = "NanoVectorDBStorage"


# Add infos
ollama_server_infos = OllamaServerInfos()

# read config.ini
config = configparser.ConfigParser()
config.read("config.ini", "utf-8")
# Redis config
redis_uri = config.get("redis", "uri", fallback=None)
if redis_uri:
    os.environ["REDIS_URI"] = redis_uri
    ollama_server_infos.KV_STORAGE = "RedisKVStorage"
    ollama_server_infos.DOC_STATUS_STORAGE = "RedisKVStorage"

# Neo4j config
neo4j_uri = config.get("neo4j", "uri", fallback=None)
neo4j_username = config.get("neo4j", "username", fallback=None)
neo4j_password = config.get("neo4j", "password", fallback=None)
if neo4j_uri:
    os.environ["NEO4J_URI"] = neo4j_uri
    os.environ["NEO4J_USERNAME"] = neo4j_username
    os.environ["NEO4J_PASSWORD"] = neo4j_password
    ollama_server_infos.GRAPH_STORAGE = "Neo4JStorage"

# Milvus config
milvus_uri = config.get("milvus", "uri", fallback=None)
milvus_user = config.get("milvus", "user", fallback=None)
milvus_password = config.get("milvus", "password", fallback=None)
milvus_db_name = config.get("milvus", "db_name", fallback=None)
if milvus_uri:
    os.environ["MILVUS_URI"] = milvus_uri
    os.environ["MILVUS_USER"] = milvus_user
    os.environ["MILVUS_PASSWORD"] = milvus_password
    os.environ["MILVUS_DB_NAME"] = milvus_db_name
    ollama_server_infos.VECTOR_STORAGE = "MilvusVectorDBStorge"

# MongoDB config
mongo_uri = config.get("mongodb", "uri", fallback=None)
mongo_database = config.get("mongodb", "MiniRAG", fallback=None)
if mongo_uri:
    os.environ["MONGO_URI"] = mongo_uri
    os.environ["MONGO_DATABASE"] = mongo_database
    ollama_server_infos.KV_STORAGE = "MongoKVStorage"
    ollama_server_infos.DOC_STATUS_STORAGE = "MongoKVStorage"


def get_default_host(binding_type: str) -> str:
    default_hosts = {
        "ollama": os.getenv("LLM_BINDING_HOST", "http://localhost:11434"),
        "lollms": os.getenv("LLM_BINDING_HOST", "http://localhost:9600"),
        "azure_openai": os.getenv("AZURE_OPENAI_ENDPOINT", "https://api.openai.com/v1"),
        "openai": os.getenv("LLM_BINDING_HOST", "https://api.openai.com/v1"),
    }
    return default_hosts.get(
        binding_type, os.getenv("LLM_BINDING_HOST", "http://localhost:11434")
    )  # fallback to ollama if unknown


def get_env_value(env_key: str, default: Any, value_type: type = str) -> Any:
    """
    Get value from environment variable with type conversion

    Args:
        env_key (str): Environment variable key
        default (Any): Default value if env variable is not set
        value_type (type): Type to convert the value to

    Returns:
        Any: Converted value from environment or default
    """
    value = os.getenv(env_key)
    if value is None:
        return default

    if isinstance(value_type, bool):
        return value.lower() in ("true", "1", "yes")
    try:
        return value_type(value)
    except ValueError:
        return default


def display_splash_screen(args: argparse.Namespace) -> None:
    """
    Display a colorful splash screen showing MiniRAG server configuration

    Args:
        args: Parsed command line arguments
    """
    # Banner
    ASCIIColors.cyan(f"""
    ╔══════════════════════════════════════════════════════════════╗
    ║                   🚀 MiniRAG Server v{__api_version__}                  ║
    ║          Fast, Lightweight RAG Server Implementation         ║
    ╚══════════════════════════════════════════════════════════════╝
    """)

    # Server Configuration
    ASCIIColors.magenta("\n📡 Server Configuration:")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.host}")
    ASCIIColors.white("    ├─ Port: ", end="")
    ASCIIColors.yellow(f"{args.port}")
    ASCIIColors.white("    ├─ SSL Enabled: ", end="")
    ASCIIColors.yellow(f"{args.ssl}")
    if args.ssl:
        ASCIIColors.white("    ├─ SSL Cert: ", end="")
        ASCIIColors.yellow(f"{args.ssl_certfile}")
        ASCIIColors.white("    └─ SSL Key: ", end="")
        ASCIIColors.yellow(f"{args.ssl_keyfile}")

    # Directory Configuration
    ASCIIColors.magenta("\n📂 Directory Configuration:")
    ASCIIColors.white("    ├─ Working Directory: ", end="")
    ASCIIColors.yellow(f"{args.working_dir}")
    ASCIIColors.white("    └─ Input Directory: ", end="")
    ASCIIColors.yellow(f"{args.input_dir}")

    # LLM Configuration
    ASCIIColors.magenta("\n🤖 LLM Configuration:")
    ASCIIColors.white("    ├─ Binding: ", end="")
    ASCIIColors.yellow(f"{args.llm_binding}")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.llm_binding_host}")
    ASCIIColors.white("    └─ Model: ", end="")
    ASCIIColors.yellow(f"{args.llm_model}")

    # Embedding Configuration
    ASCIIColors.magenta("\n📊 Embedding Configuration:")
    ASCIIColors.white("    ├─ Binding: ", end="")
    ASCIIColors.yellow(f"{args.embedding_binding}")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.embedding_binding_host}")
    ASCIIColors.white("    ├─ Model: ", end="")
    ASCIIColors.yellow(f"{args.embedding_model}")
    ASCIIColors.white("    └─ Dimensions: ", end="")
    ASCIIColors.yellow(f"{args.embedding_dim}")

    # RAG Configuration
    ASCIIColors.magenta("\n⚙️ RAG Configuration:")
    ASCIIColors.white("    ├─ Max Async Operations: ", end="")
    ASCIIColors.yellow(f"{args.max_async}")
    ASCIIColors.white("    ├─ Max Tokens: ", end="")
    ASCIIColors.yellow(f"{args.max_tokens}")
    ASCIIColors.white("    ├─ Max Embed Tokens: ", end="")
    ASCIIColors.yellow(f"{args.max_embed_tokens}")
    ASCIIColors.white("    ├─ Chunk Size: ", end="")
    ASCIIColors.yellow(f"{args.chunk_size}")
    ASCIIColors.white("    ├─ Chunk Overlap Size: ", end="")
    ASCIIColors.yellow(f"{args.chunk_overlap_size}")
    ASCIIColors.white("    ├─ History Turns: ", end="")
    ASCIIColors.yellow(f"{args.history_turns}")
    ASCIIColors.white("    ├─ Cosine Threshold: ", end="")
    ASCIIColors.yellow(f"{args.cosine_threshold}")
    ASCIIColors.white("    └─ Top-K: ", end="")
    ASCIIColors.yellow(f"{args.top_k}")

    # System Configuration
    ASCIIColors.magenta("\n🛠️ System Configuration:")
    ASCIIColors.white("    ├─ Ollama Emulating Model: ", end="")
    ASCIIColors.yellow(f"{ollama_server_infos.LIGHTRAG_MODEL}")
    ASCIIColors.white("    ├─ Log Level: ", end="")
    ASCIIColors.yellow(f"{args.log_level}")
    ASCIIColors.white("    ├─ Timeout: ", end="")
    ASCIIColors.yellow(f"{args.timeout if args.timeout else 'None (infinite)'}")
    ASCIIColors.white("    └─ API Key: ", end="")
    ASCIIColors.yellow("Set" if args.key else "Not Set")

    # Server Status
    ASCIIColors.green("\n✨ Server starting up...\n")

    # Server Access Information
    protocol = "https" if args.ssl else "http"
    if args.host == "0.0.0.0":
        ASCIIColors.magenta("\n🌐 Server Access Information:")
        ASCIIColors.white("    ├─ Local Access: ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}")
        ASCIIColors.white("    ├─ Remote Access: ", end="")
        ASCIIColors.yellow(f"{protocol}://<your-ip-address>:{args.port}")
        ASCIIColors.white("    ├─ API Documentation (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/docs")
        ASCIIColors.white("    └─ Alternative Documentation (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/redoc")

        ASCIIColors.yellow("\n📝 Note:")
        ASCIIColors.white("""    Since the server is running on 0.0.0.0:
    - Use 'localhost' or '127.0.0.1' for local access
    - Use your machine's IP address for remote access
    - To find your IP address:
      • Windows: Run 'ipconfig' in terminal
      • Linux/Mac: Run 'ifconfig' or 'ip addr' in terminal
    """)
    else:
        base_url = f"{protocol}://{args.host}:{args.port}"
        ASCIIColors.magenta("\n🌐 Server Access Information:")
        ASCIIColors.white("    ├─ Base URL: ", end="")
        ASCIIColors.yellow(f"{base_url}")
        ASCIIColors.white("    ├─ API Documentation: ", end="")
        ASCIIColors.yellow(f"{base_url}/docs")
        ASCIIColors.white("    └─ Alternative Documentation: ", end="")
        ASCIIColors.yellow(f"{base_url}/redoc")

    # Usage Examples
    ASCIIColors.magenta("\n📚 Quick Start Guide:")
    ASCIIColors.cyan("""
    1. Access the Swagger UI:
       Open your browser and navigate to the API documentation URL above

    2. API Authentication:""")
    if args.key:
        ASCIIColors.cyan("""       Add the following header to your requests:
       X-API-Key: <your-api-key>
    """)
    else:
        ASCIIColors.cyan("       No authentication required\n")

    ASCIIColors.cyan("""    3. Basic Operations:
       - POST /upload_document: Upload new documents to RAG
       - POST /query: Query your document collection
       - GET /collections: List available collections

    4. Monitor the server:
       - Check server logs for detailed operation information
       - Use healthcheck endpoint: GET /health
    """)

    # Security Notice
    if args.key:
        ASCIIColors.yellow("\n⚠️  Security Notice:")
        ASCIIColors.white("""    API Key authentication is enabled.
    Make sure to include the X-API-Key header in all your requests.
    """)

    ASCIIColors.green("Server is ready to accept connections! 🚀\n")

    # Ensure splash output flush to system log
    sys.stdout.flush()


def parse_args() -> argparse.Namespace:
    """
    Parse command line arguments with environment variable fallback

    Returns:
        argparse.Namespace: Parsed arguments
    """

    parser = argparse.ArgumentParser(
        description="MiniRAG FastAPI Server with separate working and input directories"
    )

    # Bindings configuration
    parser.add_argument(
        "--llm-binding",
        default=get_env_value("LLM_BINDING", "ollama"),
        help="LLM binding to be used. Supported: lollms, ollama, openai (default: from env or ollama)",
    )
    parser.add_argument(
        "--embedding-binding",
        default=get_env_value("EMBEDDING_BINDING", "ollama"),
        help="Embedding binding to be used. Supported: lollms, ollama, openai (default: from env or ollama)",
    )

    # Server configuration
    parser.add_argument(
        "--host",
        default=get_env_value("HOST", "0.0.0.0"),
        help="Server host (default: from env or 0.0.0.0)",
    )
    parser.add_argument(
        "--port",
        type=int,
        default=get_env_value("PORT", 9721, int),
        help="Server port (default: from env or 9721)",
    )

    # Directory configuration
    parser.add_argument(
        "--working-dir",
        default=get_env_value("WORKING_DIR", "./rag_storage"),
        help="Working directory for RAG storage (default: from env or ./rag_storage)",
    )
    parser.add_argument(
        "--input-dir",
        default=get_env_value("INPUT_DIR", "./inputs"),
        help="Directory containing input documents (default: from env or ./inputs)",
    )

    # LLM Model configuration
    parser.add_argument(
        "--llm-binding-host",
        default=get_env_value("LLM_BINDING_HOST", None),
        help="LLM server host URL. If not provided, defaults based on llm-binding:\n"
        + "- ollama: http://localhost:11434\n"
        + "- lollms: http://localhost:9600\n"
        + "- openai: https://api.openai.com/v1",
    )

    default_llm_api_key = get_env_value("LLM_BINDING_API_KEY", None)

    parser.add_argument(
        "--llm-binding-api-key",
        default=default_llm_api_key,
        help="llm server API key (default: from env or empty string)",
    )

    parser.add_argument(
        "--llm-model",
        default=get_env_value("LLM_MODEL", "mistral-nemo:latest"),
        help="LLM model name (default: from env or mistral-nemo:latest)",
    )

    # Embedding model configuration
    parser.add_argument(
        "--embedding-binding-host",
        default=get_env_value("EMBEDDING_BINDING_HOST", None),
        help="Embedding server host URL. If not provided, defaults based on embedding-binding:\n"
        + "- ollama: http://localhost:11434\n"
        + "- lollms: http://localhost:9600\n"
        + "- openai: https://api.openai.com/v1",
    )

    default_embedding_api_key = get_env_value("EMBEDDING_BINDING_API_KEY", "")
    parser.add_argument(
        "--embedding-binding-api-key",
        default=default_embedding_api_key,
        help="embedding server API key (default: from env or empty string)",
    )

    parser.add_argument(
        "--embedding-model",
        default=get_env_value("EMBEDDING_MODEL", "bge-m3:latest"),
        help="Embedding model name (default: from env or bge-m3:latest)",
    )

    parser.add_argument(
        "--chunk_size",
        default=get_env_value("CHUNK_SIZE", 1200),
        help="chunk chunk size default 1200",
    )

    parser.add_argument(
        "--chunk_overlap_size",
        default=get_env_value("CHUNK_OVERLAP_SIZE", 100),
        help="chunk overlap size default 100",
    )

    parser.add_argument(
        "--tiktoken-model-name",
        default=get_env_value("TIKTOKEN_MODEL_NAME", "cl100k_base"),
        help="Tokenizer model name for text chunking (default: from env or cl100k_base)",
    )

    def timeout_type(value):
        if value is None or value == "None":
            return None
        return int(value)

    parser.add_argument(
        "--timeout",
        default=get_env_value("TIMEOUT", None, timeout_type),
        type=timeout_type,
        help="Timeout in seconds (useful when using slow AI). Use None for infinite timeout",
    )

    # RAG configuration
    parser.add_argument(
        "--max-async",
        type=int,
        default=get_env_value("MAX_ASYNC", 4, int),
        help="Maximum async operations (default: from env or 4)",
    )
    parser.add_argument(
        "--max-tokens",
        type=int,
        default=get_env_value("MAX_TOKENS", 32768, int),
        help="Maximum token size (default: from env or 32768)",
    )
    parser.add_argument(
        "--embedding-dim",
        type=int,
        default=get_env_value("EMBEDDING_DIM", 1024, int),
        help="Embedding dimensions (default: from env or 1024)",
    )
    parser.add_argument(
        "--max-embed-tokens",
        type=int,
        default=get_env_value("MAX_EMBED_TOKENS", 8192, int),
        help="Maximum embedding token size (default: from env or 8192)",
    )

    # Logging configuration
    parser.add_argument(
        "--log-level",
        default=get_env_value("LOG_LEVEL", "INFO"),
        choices=["DEBUG", "INFO", "WARNING", "ERROR", "CRITICAL"],
        help="Logging level (default: from env or INFO)",
    )

    parser.add_argument(
        "--key",
        type=str,
        default=get_env_value("LIGHTRAG_API_KEY", None),
        help="API key for authentication. This protects minirag server against unauthorized access",
    )

    # Optional https parameters
    parser.add_argument(
        "--ssl",
        action="store_true",
        default=get_env_value("SSL", False, bool),
        help="Enable HTTPS (default: from env or False)",
    )
    parser.add_argument(
        "--ssl-certfile",
        default=get_env_value("SSL_CERTFILE", None),
        help="Path to SSL certificate file (required if --ssl is enabled)",
    )
    parser.add_argument(
        "--ssl-keyfile",
        default=get_env_value("SSL_KEYFILE", None),
        help="Path to SSL private key file (required if --ssl is enabled)",
    )
    parser.add_argument(
        "--auto-scan-at-startup",
        action="store_true",
        default=False,
        help="Enable automatic scanning when the program starts",
    )

    parser.add_argument(
        "--history-turns",
        type=int,
        default=get_env_value("HISTORY_TURNS", 3, int),
        help="Number of conversation history turns to include (default: from env or 3)",
    )

    # Search parameters
    parser.add_argument(
        "--top-k",
        type=int,
        default=get_env_value("TOP_K", 50, int),
        help="Number of most similar results to return (default: from env or 50)",
    )
    parser.add_argument(
        "--cosine-threshold",
        type=float,
        default=get_env_value("COSINE_THRESHOLD", 0.4, float),
        help="Cosine similarity threshold (default: from env or 0.4)",
    )

    parser.add_argument(
        "--simulated-model-name",
        type=str,
        default=get_env_value(
            "SIMULATED_MODEL_NAME", ollama_server_infos.LIGHTRAG_MODEL
        ),
        help="Number of conversation history turns to include (default: from env or 3)",
    )

    args = parser.parse_args()

    ollama_server_infos.LIGHTRAG_MODEL = args.simulated_model_name

    return args


class DocumentManager:
    """Handles document operations and tracking"""

    def __init__(
        self,
        input_dir: str,
        supported_extensions: tuple = (".txt", ".md", ".pdf", ".docx", ".pptx"),
    ):
        self.input_dir = Path(input_dir)
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            for file_path in self.input_dir.rglob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def scan_directory(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            for file_path in self.input_dir.rglob(f"*{ext}"):
                new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        """Mark a file as indexed"""
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        """Check if file type is supported"""
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


# Pydantic models
class SearchMode(str, Enum):
    light = "light"
    naive = "naive"
    mini = "mini"


class OllamaMessage(BaseModel):
    role: str
    content: str
    images: Optional[List[str]] = None


class OllamaChatRequest(BaseModel):
    model: str = ollama_server_infos.LIGHTRAG_MODEL
    messages: List[OllamaMessage]
    stream: bool = True  # Default to streaming mode
    options: Optional[Dict[str, Any]] = None
    system: Optional[str] = None


class OllamaChatResponse(BaseModel):
    model: str
    created_at: str
    message: OllamaMessage
    done: bool


class OllamaGenerateRequest(BaseModel):
    model: str = ollama_server_infos.LIGHTRAG_MODEL
    prompt: str
    system: Optional[str] = None
    stream: bool = False
    options: Optional[Dict[str, Any]] = None


class OllamaGenerateResponse(BaseModel):
    model: str
    created_at: str
    response: str
    done: bool
    context: Optional[List[int]]
    total_duration: Optional[int]
    load_duration: Optional[int]
    prompt_eval_count: Optional[int]
    prompt_eval_duration: Optional[int]
    eval_count: Optional[int]
    eval_duration: Optional[int]


class OllamaVersionResponse(BaseModel):
    version: str


class OllamaModelDetails(BaseModel):
    parent_model: str
    format: str
    family: str
    families: List[str]
    parameter_size: str
    quantization_level: str


class OllamaModel(BaseModel):
    name: str
    model: str
    size: int
    digest: str
    modified_at: str
    details: OllamaModelDetails


class OllamaTagResponse(BaseModel):
    models: List[OllamaModel]


class QueryRequest(BaseModel):
    query: str
    mode: SearchMode = SearchMode.light
    stream: bool = False
    only_need_context: bool = False


class QueryResponse(BaseModel):
    response: str


class InsertTextRequest(BaseModel):
    text: str
    description: Optional[str] = None


class InsertResponse(BaseModel):
    status: str
    message: str
    document_count: int


def get_api_key_dependency(api_key: Optional[str]):
    if not api_key:
        # If no API key is configured, return a dummy dependency that always succeeds
        async def no_auth():
            return None

        return no_auth

    # If API key is configured, use proper authentication
    api_key_header = APIKeyHeader(name="X-API-Key", auto_error=False)

    async def api_key_auth(api_key_header_value: str | None = Security(api_key_header)):
        if not api_key_header_value:
            raise HTTPException(
                status_code=HTTP_403_FORBIDDEN, detail="API Key required"
            )
        if api_key_header_value != api_key:
            raise HTTPException(
                status_code=HTTP_403_FORBIDDEN, detail="Invalid API Key"
            )
        return api_key_header_value

    return api_key_auth


def create_app(args):
    # Verify that bindings are correctly setup
    if args.llm_binding not in [
        "lollms",
        "ollama",
        "openai",
        "openai-ollama",
        "azure_openai",
    ]:
        raise Exception("llm binding not supported")

    if args.embedding_binding not in ["lollms", "ollama", "openai", "azure_openai"]:
        raise Exception("embedding binding not supported")

    # Set default hosts if not provided
    if args.llm_binding_host is None:
        args.llm_binding_host = get_default_host(args.llm_binding)

    if args.embedding_binding_host is None:
        args.embedding_binding_host = get_default_host(args.embedding_binding)

    # Add SSL validation
    if args.ssl:
        if not args.ssl_certfile or not args.ssl_keyfile:
            raise Exception(
                "SSL certificate and key files must be provided when SSL is enabled"
            )
        if not os.path.exists(args.ssl_certfile):
            raise Exception(f"SSL certificate file not found: {args.ssl_certfile}")
        if not os.path.exists(args.ssl_keyfile):
            raise Exception(f"SSL key file not found: {args.ssl_keyfile}")

    # Setup logging
    logging.basicConfig(
        format="%(levelname)s:%(message)s", level=getattr(logging, args.log_level)
    )

    # Check if API key is provided either through env var or args
    api_key = os.getenv("LIGHTRAG_API_KEY") or args.key

    # Initialize document manager
    doc_manager = DocumentManager(args.input_dir)

    @asynccontextmanager
    async def lifespan(app: FastAPI):
        """Lifespan context manager for startup and shutdown events"""
        # Startup logic
        if args.auto_scan_at_startup:
            try:
                new_files = doc_manager.scan_directory_for_new_files()
                for file_path in new_files:
                    try:
                        await index_file(file_path)
                    except Exception as e:
                        trace_exception(e)
                        logging.error(f"Error indexing file {file_path}: {str(e)}")

                ASCIIColors.info(
                    f"Indexed {len(new_files)} documents from {args.input_dir}"
                )
            except Exception as e:
                logging.error(f"Error during startup indexing: {str(e)}")
        yield
        # Cleanup logic (if needed)
        pass

    # Initialize FastAPI
    app = FastAPI(
        title="MiniRAG API",
        description="API for querying text using MiniRAG with separate storage and input directories"
        + "(With authentication)"
        if api_key
        else "",
        version=__api_version__,
        openapi_tags=[{"name": "api"}],
        lifespan=lifespan,
    )

    # Add CORS middleware
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # Create the optional API key dependency
    optional_api_key = get_api_key_dependency(api_key)

    # Create working directory if it doesn't exist
    Path(args.working_dir).mkdir(parents=True, exist_ok=True)
    if args.llm_binding == "lollms" or args.embedding_binding == "lollms":
        from minirag.llm.lollms import lollms_model_complete, lollms_embed
    if args.llm_binding == "ollama" or args.embedding_binding == "ollama":
        from minirag.llm.ollama import ollama_model_complete, ollama_embed
    if args.llm_binding == "openai" or args.embedding_binding == "openai":
        from minirag.llm.openai import openai_complete_if_cache, openai_embed
    if args.llm_binding == "azure_openai" or args.embedding_binding == "azure_openai":
        from minirag.llm.azure_openai import (
            azure_openai_complete_if_cache,
            azure_openai_embed,
        )
    if args.llm_binding_host == "openai-ollama" or args.embedding_binding == "ollama":
        from minirag.llm.openai import openai_complete_if_cache
        from minirag.llm.ollama import ollama_embed

    async def openai_alike_model_complete(
        prompt,
        system_prompt=None,
        history_messages=[],
        keyword_extraction=False,
        **kwargs,
    ) -> str:
        return await openai_complete_if_cache(
            args.llm_model,
            prompt,
            system_prompt=system_prompt,
            history_messages=history_messages,
            base_url=args.llm_binding_host,
            api_key=args.llm_binding_api_key,
            **kwargs,
        )

    async def azure_openai_model_complete(
        prompt,
        system_prompt=None,
        history_messages=[],
        keyword_extraction=False,
        **kwargs,
    ) -> str:
        return await azure_openai_complete_if_cache(
            args.llm_model,
            prompt,
            system_prompt=system_prompt,
            history_messages=history_messages,
            base_url=args.llm_binding_host,
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-08-01-preview"),
            **kwargs,
        )

    embedding_func = EmbeddingFunc(
        embedding_dim=args.embedding_dim,
        max_token_size=args.max_embed_tokens,
        func=lambda texts: lollms_embed(
            texts,
            embed_model=args.embedding_model,
            host=args.embedding_binding_host,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "lollms"
        else ollama_embed(
            texts,
            embed_model=args.embedding_model,
            host=args.embedding_binding_host,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "ollama"
        else azure_openai_embed(
            texts,
            model=args.embedding_model,  # no host is used for openai,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "azure_openai"
        else openai_embed(
            texts,
            model=args.embedding_model,  # no host is used for openai,
            api_key=args.embedding_binding_api_key,
        ),
    )

    # Initialize RAG
    if args.llm_binding in ["lollms", "ollama", "openai-ollama"]:
        rag = MiniRAG(
            working_dir=args.working_dir,
            llm_model_func=lollms_model_complete
            if args.llm_binding == "lollms"
            else ollama_model_complete
            if args.llm_binding == "ollama"
            else openai_alike_model_complete,
            llm_model_name=args.llm_model,
            llm_model_max_async=args.max_async,
            llm_model_max_token_size=args.max_tokens,
            chunk_token_size=int(args.chunk_size),
            chunk_overlap_token_size=int(args.chunk_overlap_size),
            tiktoken_model_name=args.tiktoken_model_name,
            llm_model_kwargs={
                "host": args.llm_binding_host,
                "timeout": args.timeout,
                "options": {"num_ctx": args.max_tokens},
                "api_key": args.llm_binding_api_key,
            }
            if args.llm_binding == "lollms" or args.llm_binding == "ollama"
            else {},
            embedding_func=embedding_func,
            kv_storage=ollama_server_infos.KV_STORAGE,
            graph_storage=ollama_server_infos.GRAPH_STORAGE,
            vector_storage=ollama_server_infos.VECTOR_STORAGE,
            doc_status_storage=ollama_server_infos.DOC_STATUS_STORAGE,
            vector_db_storage_cls_kwargs={
                "cosine_better_than_threshold": args.cosine_threshold
            },
        )
    else:
        rag = MiniRAG(
            working_dir=args.working_dir,
            llm_model_func=azure_openai_model_complete
            if args.llm_binding == "azure_openai"
            else openai_alike_model_complete,
            chunk_token_size=int(args.chunk_size),
            chunk_overlap_token_size=int(args.chunk_overlap_size),
            tiktoken_model_name=args.tiktoken_model_name,
            llm_model_kwargs={
                "timeout": args.timeout,
            },
            llm_model_name=args.llm_model,
            llm_model_max_async=args.max_async,
            llm_model_max_token_size=args.max_tokens,
            embedding_func=embedding_func,
            kv_storage=ollama_server_infos.KV_STORAGE,
            graph_storage=ollama_server_infos.GRAPH_STORAGE,
            vector_storage=ollama_server_infos.VECTOR_STORAGE,
            doc_status_storage=ollama_server_infos.DOC_STATUS_STORAGE,
            vector_db_storage_cls_kwargs={
                "cosine_better_than_threshold": args.cosine_threshold
            },
        )

    async def index_file(file_path: Union[str, Path]) -> None:
        """Index all files inside the folder with support for multiple file formats

        Args:
            file_path: Path to the file to be indexed (str or Path object)

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        if not pm.is_installed("aiofiles"):
            pm.install("aiofiles")

        # Convert to Path object if string
        file_path = Path(file_path)

        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        content = ""
        # Get file extension in lowercase
        ext = file_path.suffix.lower()

        if ext in [".txt", ".md"]:
            # Text files handling
            async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                content = await f.read()
        elif ext == ".pdf":
            if not pm.is_installed("pypdf2"):
                pm.install("pypdf2")
            from PyPDF2 import PdfReader

            # PDF handling
            reader = PdfReader(str(file_path))
            content = ""
            for page in reader.pages:
                content += page.extract_text() + "\n"
        elif ext == ".docx":
            if not pm.is_installed("python-docx"):
                pm.install("python-docx")
            from docx import Document

            # Word document handling
            doc = Document(file_path)
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif ext == ".pptx":
            if not pm.is_installed("pptx"):
                pm.install("pptx")
            from pptx import Presentation  # type: ignore

            # PowerPoint handling
            prs = Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        else:
            raise ValueError(f"Unsupported file format: {ext}")

        # Insert content into RAG system
        if content:
            await rag.ainsert(content)
            doc_manager.mark_as_indexed(file_path)
            logging.info(f"Successfully indexed file: {file_path}")
        else:
            logging.warning(f"No content extracted from file: {file_path}")

    @app.post("/documents/scan", dependencies=[Depends(optional_api_key)])
    async def scan_for_new_documents():
        """Trigger the scanning process"""
        global scan_progress

        try:
            with progress_lock:
                if scan_progress["is_scanning"]:
                    return {"status": "already_scanning"}

                scan_progress["is_scanning"] = True
                scan_progress["indexed_count"] = 0
                scan_progress["progress"] = 0

            new_files = doc_manager.scan_directory_for_new_files()
            scan_progress["total_files"] = len(new_files)

            for file_path in new_files:
                try:
                    with progress_lock:
                        scan_progress["current_file"] = os.path.basename(file_path)

                    await index_file(file_path)

                    with progress_lock:
                        scan_progress["indexed_count"] += 1
                        scan_progress["progress"] = (
                            scan_progress["indexed_count"]
                            / scan_progress["total_files"]
                        ) * 100

                except Exception as e:
                    logging.error(f"Error indexing file {file_path}: {str(e)}")

            return {
                "status": "success",
                "indexed_count": scan_progress["indexed_count"],
                "total_documents": len(doc_manager.indexed_files),
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))
        finally:
            with progress_lock:
                scan_progress["is_scanning"] = False

    @app.get("/documents/scan-progress")
    async def get_scan_progress():
        """Get the current scanning progress"""
        with progress_lock:
            return scan_progress

    @app.post("/documents/upload", dependencies=[Depends(optional_api_key)])
    async def upload_to_input_dir(file: UploadFile = File(...)):
        """
        Endpoint for uploading a file to the input directory and indexing it.

        This API endpoint accepts a file through an HTTP POST request, checks if the
        uploaded file is of a supported type, saves it in the specified input directory,
        indexes it for retrieval, and returns a success status with relevant details.

        Parameters:
            file (UploadFile): The file to be uploaded. It must have an allowed extension as per
                               `doc_manager.supported_extensions`.

        Returns:
            dict: A dictionary containing the upload status ("success"),
                  a message detailing the operation result, and
                  the total number of indexed documents.

        Raises:
            HTTPException: If the file type is not supported, it raises a 400 Bad Request error.
                           If any other exception occurs during the file handling or indexing,
                           it raises a 500 Internal Server Error with details about the exception.
        """
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            file_path = doc_manager.input_dir / file.filename
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Immediately index the uploaded file
            await index_file(file_path)

            return {
                "status": "success",
                "message": f"File uploaded and indexed: {file.filename}",
                "total_documents": len(doc_manager.indexed_files),
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/query", response_model=QueryResponse, dependencies=[Depends(optional_api_key)]
    )
    async def query_text(request: QueryRequest):
        """
        Handle a POST request at the /query endpoint to process user queries using RAG capabilities.

        Parameters:
            request (QueryRequest): A Pydantic model containing the following fields:
                - query (str): The text of the user's query.
                - mode (ModeEnum): Optional. Specifies the mode of retrieval augmentation.
                - stream (bool): Optional. Determines if the response should be streamed.
                - only_need_context (bool): Optional. If true, returns only the context without further processing.

        Returns:
            QueryResponse: A Pydantic model containing the result of the query processing.
                           If a string is returned (e.g., cache hit), it's directly returned.
                           Otherwise, an async generator may be used to build the response.

        Raises:
            HTTPException: Raised when an error occurs during the request handling process,
                           with status code 500 and detail containing the exception message.
        """
        try:
            response = await rag.aquery(
                request.query,
                param=QueryParam(
                    mode=request.mode,
                    stream=request.stream,
                    only_need_context=request.only_need_context,
                    top_k=args.top_k,
                ),
            )

            # If response is a string (e.g. cache hit), return directly
            if isinstance(response, str):
                return QueryResponse(response=response)

            # If it's an async generator, decide whether to stream based on stream parameter
            if request.stream:
                result = ""
                async for chunk in response:
                    result += chunk
                return QueryResponse(response=result)
            else:
                result = ""
                async for chunk in response:
                    result += chunk
                return QueryResponse(response=result)
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @app.post("/query/stream", dependencies=[Depends(optional_api_key)])
    async def query_text_stream(request: QueryRequest):
        """
        This endpoint performs a retrieval-augmented generation (RAG) query and streams the response.

        Args:
            request (QueryRequest): The request object containing the query parameters.
            optional_api_key (Optional[str], optional): An optional API key for authentication. Defaults to None.

        Returns:
            StreamingResponse: A streaming response containing the RAG query results.
        """
        try:
            response = await rag.aquery(  # Use aquery instead of query, and add await
                request.query,
                param=QueryParam(
                    mode=request.mode,
                    stream=True,
                    only_need_context=request.only_need_context,
                    top_k=args.top_k,
                ),
            )

            from fastapi.responses import StreamingResponse

            async def stream_generator():
                if isinstance(response, str):
                    # If it's a string, send it all at once
                    yield f"{json.dumps({'response': response})}\n"
                else:
                    # If it's an async generator, send chunks one by one
                    try:
                        async for chunk in response:
                            if chunk:  # Only send non-empty content
                                yield f"{json.dumps({'response': chunk})}\n"
                    except Exception as e:
                        logging.error(f"Streaming error: {str(e)}")
                        yield f"{json.dumps({'error': str(e)})}\n"

            return StreamingResponse(
                stream_generator(),
                media_type="application/x-ndjson",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "Content-Type": "application/x-ndjson",
                    "Access-Control-Allow-Origin": "*",
                    "Access-Control-Allow-Methods": "POST, OPTIONS",
                    "Access-Control-Allow-Headers": "Content-Type",
                    "X-Accel-Buffering": "no",  # Disable Nginx buffering
                },
            )
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/text",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_text(request: InsertTextRequest):
        """
        Insert text into the Retrieval-Augmented Generation (RAG) system.

        This endpoint allows you to insert text data into the RAG system for later retrieval and use in generating responses.

        Args:
            request (InsertTextRequest): The request body containing the text to be inserted.

        Returns:
            InsertResponse: A response object containing the status of the operation, a message, and the number of documents inserted.
        """
        import time
        start_time = time.time()
        
        try:
            logging.info(f"开始插入文本 (长度: {len(request.text)} 字符)")
            
            # 插入到MiniRAG
            insert_start = time.time()
            await rag.ainsert(request.text)
            insert_time = time.time() - insert_start
            
            total_time = time.time() - start_time
            logging.info(f"文本插入成功 (MiniRAG耗时: {insert_time:.3f}s, 总耗时: {total_time:.3f}s)")
            
            return InsertResponse(
                status="success",
                message="Text successfully inserted",
                document_count=1,
            )
        except Exception as e:
            total_time = time.time() - start_time
            logging.error(f"文本插入失败 (耗时: {total_time:.3f}s): {e}")
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/file",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_file(file: UploadFile = File(...), description: str = Form(None)):
        """Insert a file directly into the RAG system

        Args:
            file: Uploaded file
            description: Optional description of the file

        Returns:
            InsertResponse: Status of the insertion operation

        Raises:
            HTTPException: For unsupported file types or processing errors
        """
        try:
            content = ""
            # Get file extension in lowercase
            ext = Path(file.filename).suffix.lower()

            if ext in [".txt", ".md"]:
                # Text files handling
                text_content = await file.read()
                content = text_content.decode("utf-8")
            elif ext == ".pdf":
                if not pm.is_installed("pypdf2"):
                    pm.install("pypdf2")
                from PyPDF2 import PdfReader
                from io import BytesIO

                # Read PDF from memory
                pdf_content = await file.read()
                pdf_file = BytesIO(pdf_content)
                reader = PdfReader(pdf_file)
                content = ""
                for page in reader.pages:
                    content += page.extract_text() + "\n"
            elif ext == ".docx":
                if not pm.is_installed("python-docx"):
                    pm.install("python-docx")
                from docx import Document
                from io import BytesIO

                # Read DOCX from memory
                docx_content = await file.read()
                docx_file = BytesIO(docx_content)
                doc = Document(docx_file)
                content = "\n".join(
                    [paragraph.text for paragraph in doc.paragraphs]
                )
            elif ext == ".pptx":
                if not pm.is_installed("pptx"):
                    pm.install("pptx")
                from pptx import Presentation  # type: ignore
                from io import BytesIO

                # Read PPTX from memory
                pptx_content = await file.read()
                pptx_file = BytesIO(pptx_content)
                prs = Presentation(pptx_file)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            else:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            # Insert content into RAG system
            if content:
                # Add description if provided
                if description:
                    content = f"{description}\n\n{content}"

                await rag.ainsert(content)
                logging.info(f"Successfully indexed file: {file.filename}")

                return InsertResponse(
                    status="success",
                    message=f"File '{file.filename}' successfully inserted",
                    document_count=1,
                )
            else:
                raise HTTPException(
                    status_code=400,
                    detail="No content could be extracted from the file",
                )

        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="File encoding not supported")
        except Exception as e:
            logging.error(f"Error processing file {file.filename}: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/batch",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_batch(files: List[UploadFile] = File(...)):
        """Process multiple files in batch mode

        Args:
            files: List of files to process

        Returns:
            InsertResponse: Status of the batch insertion operation

        Raises:
            HTTPException: For processing errors
        """
        try:
            inserted_count = 0
            failed_files = []

            for file in files:
                try:
                    content = ""
                    ext = Path(file.filename).suffix.lower()

                    if ext in [".txt", ".md"]:
                        text_content = await file.read()
                        content = text_content.decode("utf-8")
                    elif ext == ".pdf":
                        if not pm.is_installed("pypdf2"):
                            pm.install("pypdf2")
                        from PyPDF2 import PdfReader
                        from io import BytesIO

                        pdf_content = await file.read()
                        pdf_file = BytesIO(pdf_content)
                        reader = PdfReader(pdf_file)
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                    elif ext == ".docx":
                        if not pm.is_installed("docx"):
                            pm.install("docx")
                        from docx import Document
                        from io import BytesIO

                        docx_content = await file.read()
                        docx_file = BytesIO(docx_content)
                        doc = Document(docx_file)
                        content = "\n".join(
                            [paragraph.text for paragraph in doc.paragraphs]
                        )
                    elif ext == ".pptx":
                        if not pm.is_installed("pptx"):
                            pm.install("pptx")
                        from pptx import Presentation  # type: ignore
                        from io import BytesIO

                        pptx_content = await file.read()
                        pptx_file = BytesIO(pptx_content)
                        prs = Presentation(pptx_file)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    content += shape.text + "\n"
                    else:
                        failed_files.append(f"{file.filename} (unsupported type)")
                        continue

                    if content:
                        await rag.ainsert(content)
                        inserted_count += 1
                        logging.info(f"Successfully indexed file: {file.filename}")
                    else:
                        failed_files.append(f"{file.filename} (no content extracted)")

                except UnicodeDecodeError:
                    failed_files.append(f"{file.filename} (encoding error)")
                except Exception as e:
                    failed_files.append(f"{file.filename} ({str(e)})")
                    logging.error(f"Error processing file {file.filename}: {str(e)}")

            # Prepare status message
            if inserted_count == len(files):
                status = "success"
                status_message = f"Successfully inserted all {inserted_count} documents"
            elif inserted_count > 0:
                status = "partial_success"
                status_message = f"Successfully inserted {inserted_count} out of {len(files)} documents"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"
            else:
                status = "failure"
                status_message = "No documents were successfully inserted"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"

            return InsertResponse(
                status=status,
                message=status_message,
                document_count=inserted_count,
            )

        except Exception as e:
            logging.error(f"Batch processing error: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    @app.delete(
        "/documents",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def clear_documents():
        """
        Clear all documents from the MiniRAG system.

        This endpoint deletes all text chunks, entities vector database, and relationships vector database,
        effectively clearing all documents from the MiniRAG system.

        Returns:
            InsertResponse: A response object containing the status, message, and the new document count (0 in this case).
        """
        try:
            # 清空存储内容而不是覆盖存储对象
            if hasattr(rag.text_chunks, 'drop'):
                rag.text_chunks.drop()
            if hasattr(rag.entities_vdb, 'drop'):
                rag.entities_vdb.drop()
            if hasattr(rag.relationships_vdb, 'drop'):
                rag.relationships_vdb.drop()
            if hasattr(rag.doc_status, 'drop'):
                rag.doc_status.drop()
            if hasattr(rag.full_docs, 'drop'):
                rag.full_docs.drop()
            
            return InsertResponse(
                status="success",
                message="All documents cleared successfully",
                document_count=0,
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    # query all graph labels
    @app.get("/graph/label/list")
    async def get_graph_labels():
        return await rag.get_graph_labels()

    # query all graph
    @app.get("/graphs")
    async def get_graphs(label: str):
        return await rag.get_graps(nodel_label=label, max_depth=100)

    # Ollama compatible API endpoints
    # -------------------------------------------------
    @app.get("/api/version")
    async def get_version():
        """Get Ollama version information"""
        return OllamaVersionResponse(version="0.5.4")

    @app.get("/api/tags")
    async def get_tags():
        """Get available models"""
        return OllamaTagResponse(
            models=[
                {
                    "name": ollama_server_infos.LIGHTRAG_MODEL,
                    "model": ollama_server_infos.LIGHTRAG_MODEL,
                    "size": ollama_server_infos.LIGHTRAG_SIZE,
                    "digest": ollama_server_infos.LIGHTRAG_DIGEST,
                    "modified_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                    "details": {
                        "parent_model": "",
                        "format": "gguf",
                        "family": ollama_server_infos.LIGHTRAG_NAME,
                        "families": [ollama_server_infos.LIGHTRAG_NAME],
                        "parameter_size": "13B",
                        "quantization_level": "Q4_0",
                    },
                }
            ]
        )

    def parse_query_mode(query: str) -> tuple[str, SearchMode]:
        """Parse query prefix to determine search mode
        Returns tuple of (cleaned_query, search_mode)
        """
        mode_map = {
            "/light ": SearchMode.light,
            "/naive ": SearchMode.naive,
            "/mini ": SearchMode.mini,
        }

        for prefix, mode in mode_map.items():
            if query.startswith(prefix):
                # After removing prefix an leading spaces
                cleaned_query = query[len(prefix) :].lstrip()
                return cleaned_query, mode

        return query, SearchMode.hybrid

    @app.post("/api/generate")
    async def generate(raw_request: Request, request: OllamaGenerateRequest):
        """Handle generate completion requests
        For compatiblity purpuse, the request is not processed by MiniRAG,
        and will be handled by underlying LLM model.
        """
        try:
            query = request.prompt
            start_time = time.time_ns()
            prompt_tokens = estimate_tokens(query)

            if request.system:
                rag.llm_model_kwargs["system_prompt"] = request.system

            if request.stream:
                from fastapi.responses import StreamingResponse

                response = await rag.llm_model_func(
                    query, stream=True, **rag.llm_model_kwargs
                )

                async def stream_generator():
                    try:
                        first_chunk_time = None
                        last_chunk_time = None
                        total_response = ""

                        # Ensure response is an async generator
                        if isinstance(response, str):
                            # If it's a string, send in two parts
                            first_chunk_time = time.time_ns()
                            last_chunk_time = first_chunk_time
                            total_response = response

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "response": response,
                                "done": False,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"

                            completion_tokens = estimate_tokens(total_response)
                            total_time = last_chunk_time - start_time
                            prompt_eval_time = first_chunk_time - start_time
                            eval_time = last_chunk_time - first_chunk_time

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "done": True,
                                "total_duration": total_time,
                                "load_duration": 0,
                                "prompt_eval_count": prompt_tokens,
                                "prompt_eval_duration": prompt_eval_time,
                                "eval_count": completion_tokens,
                                "eval_duration": eval_time,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"
                        else:
                            async for chunk in response:
                                if chunk:
                                    if first_chunk_time is None:
                                        first_chunk_time = time.time_ns()

                                    last_chunk_time = time.time_ns()

                                    total_response += chunk
                                    data = {
                                        "model": ollama_server_infos.LIGHTRAG_MODEL,
                                        "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                        "response": chunk,
                                        "done": False,
                                    }
                                    yield f"{json.dumps(data, ensure_ascii=False)}\n"

                            completion_tokens = estimate_tokens(total_response)
                            total_time = last_chunk_time - start_time
                            prompt_eval_time = first_chunk_time - start_time
                            eval_time = last_chunk_time - first_chunk_time

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "done": True,
                                "total_duration": total_time,
                                "load_duration": 0,
                                "prompt_eval_count": prompt_tokens,
                                "prompt_eval_duration": prompt_eval_time,
                                "eval_count": completion_tokens,
                                "eval_duration": eval_time,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"
                            return

                    except Exception as e:
                        logging.error(f"Error in stream_generator: {str(e)}")
                        raise

                return StreamingResponse(
                    stream_generator(),
                    media_type="application/x-ndjson",
                    headers={
                        "Cache-Control": "no-cache",
                        "Connection": "keep-alive",
                        "Content-Type": "application/x-ndjson",
                        "Access-Control-Allow-Origin": "*",
                        "Access-Control-Allow-Methods": "POST, OPTIONS",
                        "Access-Control-Allow-Headers": "Content-Type",
                    },
                )
            else:
                first_chunk_time = time.time_ns()
                response_text = await rag.llm_model_func(
                    query, stream=False, **rag.llm_model_kwargs
                )
                last_chunk_time = time.time_ns()

                if not response_text:
                    response_text = "No response generated"

                completion_tokens = estimate_tokens(str(response_text))
                total_time = last_chunk_time - start_time
                prompt_eval_time = first_chunk_time - start_time
                eval_time = last_chunk_time - first_chunk_time

                return {
                    "model": ollama_server_infos.LIGHTRAG_MODEL,
                    "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                    "response": str(response_text),
                    "done": True,
                    "total_duration": total_time,
                    "load_duration": 0,
                    "prompt_eval_count": prompt_tokens,
                    "prompt_eval_duration": prompt_eval_time,
                    "eval_count": completion_tokens,
                    "eval_duration": eval_time,
                }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @app.post("/api/chat")
    async def chat(raw_request: Request, request: OllamaChatRequest):
        """Process chat completion requests.
        Routes user queries through MiniRAG by selecting query mode based on prefix indicators.
        Detects and forwards OpenWebUI session-related requests (for meta data generation task) directly to LLM.
        """
        try:
            # Get all messages
            messages = request.messages
            if not messages:
                raise HTTPException(status_code=400, detail="No messages provided")

            # Get the last message as query and previous messages as history
            query = messages[-1].content
            # Convert OllamaMessage objects to dictionaries
            conversation_history = [
                {"role": msg.role, "content": msg.content} for msg in messages[:-1]
            ]

            # Check for query prefix
            cleaned_query, mode = parse_query_mode(query)

            start_time = time.time_ns()
            prompt_tokens = estimate_tokens(cleaned_query)

            param_dict = {
                "mode": mode,
                "stream": request.stream,
                "only_need_context": False,
                "conversation_history": conversation_history,
                "top_k": args.top_k,
            }

            if args.history_turns is not None:
                param_dict["history_turns"] = args.history_turns

            query_param = QueryParam(**param_dict)

            if request.stream:
                from fastapi.responses import StreamingResponse

                response = await rag.aquery(  # Need await to get async generator
                    cleaned_query, param=query_param
                )

                async def stream_generator():
                    try:
                        first_chunk_time = None
                        last_chunk_time = None
                        total_response = ""

                        # Ensure response is an async generator
                        if isinstance(response, str):
                            # If it's a string, send in two parts
                            first_chunk_time = time.time_ns()
                            last_chunk_time = first_chunk_time
                            total_response = response

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "message": {
                                    "role": "assistant",
                                    "content": response,
                                    "images": None,
                                },
                                "done": False,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"

                            completion_tokens = estimate_tokens(total_response)
                            total_time = last_chunk_time - start_time
                            prompt_eval_time = first_chunk_time - start_time
                            eval_time = last_chunk_time - first_chunk_time

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "done": True,
                                "total_duration": total_time,
                                "load_duration": 0,
                                "prompt_eval_count": prompt_tokens,
                                "prompt_eval_duration": prompt_eval_time,
                                "eval_count": completion_tokens,
                                "eval_duration": eval_time,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"
                        else:
                            async for chunk in response:
                                if chunk:
                                    if first_chunk_time is None:
                                        first_chunk_time = time.time_ns()

                                    last_chunk_time = time.time_ns()

                                    total_response += chunk
                                    data = {
                                        "model": ollama_server_infos.LIGHTRAG_MODEL,
                                        "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                        "message": {
                                            "role": "assistant",
                                            "content": chunk,
                                            "images": None,
                                        },
                                        "done": False,
                                    }
                                    yield f"{json.dumps(data, ensure_ascii=False)}\n"

                            completion_tokens = estimate_tokens(total_response)
                            total_time = last_chunk_time - start_time
                            prompt_eval_time = first_chunk_time - start_time
                            eval_time = last_chunk_time - first_chunk_time

                            data = {
                                "model": ollama_server_infos.LIGHTRAG_MODEL,
                                "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                                "done": True,
                                "total_duration": total_time,
                                "load_duration": 0,
                                "prompt_eval_count": prompt_tokens,
                                "prompt_eval_duration": prompt_eval_time,
                                "eval_count": completion_tokens,
                                "eval_duration": eval_time,
                            }
                            yield f"{json.dumps(data, ensure_ascii=False)}\n"
                            return  # Ensure the generator ends immediately after sending the completion marker
                    except Exception as e:
                        logging.error(f"Error in stream_generator: {str(e)}")
                        raise

                return StreamingResponse(
                    stream_generator(),
                    media_type="application/x-ndjson",
                    headers={
                        "Cache-Control": "no-cache",
                        "Connection": "keep-alive",
                        "Content-Type": "application/x-ndjson",
                        "Access-Control-Allow-Origin": "*",
                        "Access-Control-Allow-Methods": "POST, OPTIONS",
                        "Access-Control-Allow-Headers": "Content-Type",
                    },
                )
            else:
                first_chunk_time = time.time_ns()

                # Determine if the request is from Open WebUI's session title and session keyword generation task
                match_result = re.search(
                    r"\n<chat_history>\nUSER:", cleaned_query, re.MULTILINE
                )
                if match_result:
                    if request.system:
                        rag.llm_model_kwargs["system_prompt"] = request.system

                    response_text = await rag.llm_model_func(
                        cleaned_query, stream=False, **rag.llm_model_kwargs
                    )
                else:
                    response_text = await rag.aquery(cleaned_query, param=query_param)

                last_chunk_time = time.time_ns()

                if not response_text:
                    response_text = "No response generated"

                completion_tokens = estimate_tokens(str(response_text))
                total_time = last_chunk_time - start_time
                prompt_eval_time = first_chunk_time - start_time
                eval_time = last_chunk_time - first_chunk_time

                return {
                    "model": ollama_server_infos.LIGHTRAG_MODEL,
                    "created_at": ollama_server_infos.LIGHTRAG_CREATED_AT,
                    "message": {
                        "role": "assistant",
                        "content": str(response_text),
                        "images": None,
                    },
                    "done": True,
                    "total_duration": total_time,
                    "load_duration": 0,
                    "prompt_eval_count": prompt_tokens,
                    "prompt_eval_duration": prompt_eval_time,
                    "eval_count": completion_tokens,
                    "eval_duration": eval_time,
                }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @app.get("/documents", dependencies=[Depends(optional_api_key)])
    async def documents():
        """Get current system status"""
        try:
            # 直接读取文档数据 from MiniRAG storage backends
            doc_status_data = {}
            if rag.doc_status:
                try:
                    doc_status_data = rag.doc_status.all_keys()
                except Exception as e:
                    logging.warning(f"获取文档状态失败: {e}")
            
            full_docs_data = {}
            if rag.full_docs:
                try:
                    full_docs_data = rag.full_docs.all_keys()
                except Exception as e:
                    logging.warning(f"获取完整文档失败: {e}")
            
            text_chunks_data = {}
            if rag.text_chunks:
                try:
                    text_chunks_data = rag.text_chunks.all_keys()
                except Exception as e:
                    logging.warning(f"获取文本块失败: {e}")
            
            documents = []
            for doc_id in doc_status_data:
                try:
                    doc_status = rag.doc_status.get(doc_id)
                    if doc_status and doc_status.get("status") == "PROCESSED":
                        doc_info = {
                            "id": doc_id,
                            "status": doc_status.get("status"),
                            "chunks_count": doc_status.get("chunks_count", 0),
                            "content_length": doc_status.get("content_length", 0),
                            "created_at": doc_status.get("created_at"),
                            "updated_at": doc_status.get("updated_at")
                        }
                        if "content_summary" in doc_status:
                            doc_info["content_summary"] = doc_status["content_summary"]
                        
                        # 获取完整文档内容
                        if doc_id in full_docs_data:
                            try:
                                full_doc = rag.full_docs.get(doc_id)
                                if full_doc:
                                    doc_info["content"] = full_doc.get("content", "")
                            except Exception as e:
                                logging.warning(f"获取文档内容失败 {doc_id}: {e}")
                        
                        documents.append(doc_info)
                except Exception as e:
                    logging.warning(f"处理文档 {doc_id} 失败: {e}")
                    continue
            
            return {
                "status": "success",
                "total_documents": len(documents),
                "documents": documents,
                "storage_info": {
                    "doc_status_count": len(doc_status_data),
                    "full_docs_count": len(full_docs_data),
                    "text_chunks_count": len(text_chunks_data)
                }
            }
        except Exception as e:
            logging.error(f"获取文档列表失败: {e}")
            return {
                "status": "error",
                "message": str(e),
                "total_documents": 0,
                "documents": []
            }

    @app.get("/health", dependencies=[Depends(optional_api_key)])
    async def get_status():
        """Get current system status"""
        files = doc_manager.scan_directory()
        return {
            "status": "healthy",
            "working_directory": str(args.working_dir),
            "input_directory": str(args.input_dir),
            "indexed_files": [str(f) for f in files],
            "indexed_files_count": len(files),
            "configuration": {
                # LLM configuration binding/host address (if applicable)/model (if applicable)
                "llm_binding": args.llm_binding,
                "llm_binding_host": args.llm_binding_host,
                "llm_model": args.llm_model,
                # embedding model configuration binding/host address (if applicable)/model (if applicable)
                "embedding_binding": args.embedding_binding,
                "embedding_binding_host": args.embedding_binding_host,
                "embedding_model": args.embedding_model,
                "max_tokens": args.max_tokens,
                "kv_storage": ollama_server_infos.KV_STORAGE,
                "doc_status_storage": ollama_server_infos.DOC_STATUS_STORAGE,
                "graph_storage": ollama_server_infos.GRAPH_STORAGE,
                "vector_storage": ollama_server_infos.VECTOR_STORAGE,
            },
        }

    # 文件上传和解析相关API
    @app.post("/upload/file", dependencies=[Depends(optional_api_key)])
    async def upload_file(
        background_tasks: BackgroundTasks,
        file: UploadFile = File(...),
        description: str = Form("")
    ):
        """上传并解析文件"""
        import time
        start_time = time.time()
        
        try:
            logging.info(f"开始处理文件: {file.filename} (大小: {file.size} bytes)")
            
            # 创建上传目录
            upload_start = time.time()
            upload_dir = Path(args.working_dir) / "uploads"
            upload_dir.mkdir(parents=True, exist_ok=True)
            upload_time = time.time() - upload_start
            logging.info(f"创建上传目录耗时: {upload_time:.3f}s")
            
            # 保存文件
            save_start = time.time()
            file_path = upload_dir / file.filename
            with open(file_path, "wb") as buffer:
                content = await file.read()
                buffer.write(content)
            save_time = time.time() - save_start
            logging.info(f"文件保存耗时: {save_time:.3f}s (大小: {len(content)} bytes)")
            
            # 后台任务：解析并索引文件
            async def process_file():
                process_start = time.time()
                try:
                    logging.info(f"开始后台处理文件: {file_path}")
                    
                    # 解析文件
                    parse_start = time.time()
                    from minirag.file_parser import parse_single_file
                    logging.info(f"开始解析文件: {file_path}")
                    parsed_result = await parse_single_file(str(file_path))
                    parse_time = time.time() - parse_start
                    
                    if 'error' in parsed_result:
                        logging.error(f"文件解析失败 {file_path}: {parsed_result['error']} (耗时: {parse_time:.3f}s)")
                        return
                    
                    logging.info(f"文件解析成功 (耗时: {parse_time:.3f}s)")
                    logging.info(f"解析内容长度: {len(parsed_result.get('content', ''))} 字符")
                    
                    # 如果有描述，添加到内容中
                    if description:
                        content = f"文件描述: {description}\n\n文件内容:\n{parsed_result.get('content', '')}"
                    else:
                        content = parsed_result.get('content', '')
                    
                    # 插入到MiniRAG
                    insert_start = time.time()
                    logging.info(f"开始插入到MiniRAG (内容长度: {len(content)} 字符)")
                    await rag.ainsert(content)
                    insert_time = time.time() - insert_start
                    logging.info(f"MiniRAG插入成功 (耗时: {insert_time:.3f}s)")
                    
                    total_process_time = time.time() - process_start
                    logging.info(f"文件处理完成: {file_path} (总耗时: {total_process_time:.3f}s)")
                    
                except Exception as e:
                    process_time = time.time() - process_start
                    logging.error(f"处理上传文件失败 {file_path} (耗时: {process_time:.3f}s): {e}")
            
            background_tasks.add_task(process_file)
            
            total_time = time.time() - start_time
            logging.info(f"文件上传响应完成: {file.filename} (总耗时: {total_time:.3f}s)")
            
            return {
                "status": "success",
                "message": f"文件 {file.filename} 上传成功，正在后台处理",
                "file_path": str(file_path),
                "file_size": len(content),
                "processing_time": {
                    "upload_dir_creation": f"{upload_time:.3f}s",
                    "file_save": f"{save_time:.3f}s",
                    "total_response_time": f"{total_time:.3f}s"
                }
            }
        except Exception as e:
            total_time = time.time() - start_time
            logging.error(f"文件上传失败: {file.filename} (耗时: {total_time:.3f}s): {str(e)}")
            raise HTTPException(status_code=500, detail=f"文件上传失败: {str(e)}")

    @app.post("/upload/image", dependencies=[Depends(optional_api_key)])
    async def upload_image(
        background_tasks: BackgroundTasks,
        image: UploadFile = File(...),
        description: str = Form("")
    ):
        """上传并解析图片（OCR）"""
        try:
            # 检查文件类型
            if not image.content_type.startswith('image/'):
                raise HTTPException(status_code=400, detail="只支持图片文件")
            
            # 创建上传目录
            upload_dir = Path(args.working_dir) / "images"
            upload_dir.mkdir(parents=True, exist_ok=True)
            
            # 保存图片
            image_path = upload_dir / image.filename
            with open(image_path, "wb") as buffer:
                content = await image.read()
                buffer.write(content)
            
            # 后台任务：多模态解析并索引图片
            async def process_image():
                try:
                    from minirag.file_parser import parse_single_file
                    # 解析图片（多模态理解或OCR）
                    parsed_result = await parse_single_file(str(image_path))
                    
                    if 'error' in parsed_result:
                        logging.error(f"图片解析失败 {image_path}: {parsed_result['error']}")
                        return
                    
                    # 构建内容（包含图片描述和多模态理解结果）
                    processing_method = parsed_result.get('processing_method', 'unknown')
                    if description:
                        content = f"图片描述: {description}\n\n图片内容({processing_method}):\n{parsed_result.get('content', '')}"
                    else:
                        content = f"图片内容({processing_method}):\n{parsed_result.get('content', '')}"
                    
                    # 插入到MiniRAG
                    await rag.ainsert(content)
                    
                    logging.info(f"图片处理成功: {image_path} (方法: {processing_method})")
                    
                except Exception as e:
                    logging.error(f"处理上传图片失败 {image_path}: {e}")
            
            background_tasks.add_task(process_image)
            
            return {
                "status": "success",
                "message": f"图片 {image.filename} 上传成功，正在后台处理",
                "image_path": str(image_path),
                "file_size": len(content)
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"图片上传失败: {str(e)}")

    @app.post("/query/image", dependencies=[Depends(optional_api_key)])
    async def query_image(
        image: UploadFile = File(...),
        question: str = Form(...),
        mode: str = Form("mini")
    ):
        """基于图片的问答"""
        try:
            # 检查文件类型
            if not image.content_type.startswith('image/'):
                raise HTTPException(status_code=400, detail="只支持图片文件")
            
            # 临时保存图片
            temp_dir = Path(args.working_dir) / "temp"
            temp_dir.mkdir(parents=True, exist_ok=True)
            
            temp_image_path = temp_dir / f"temp_{int(time.time())}_{image.filename}"
            with open(temp_image_path, "wb") as buffer:
                content = await image.read()
                buffer.write(content)
            
            try:
                # 解析图片内容
                from minirag.file_parser import parse_single_file
                parsed_result = await parse_single_file(str(temp_image_path))
                
                if 'error' in parsed_result:
                    raise HTTPException(status_code=500, detail=f"图片解析失败: {parsed_result['error']}")
                
                # 构建查询（包含图片内容描述）
                image_context = f"图片内容: {parsed_result.get('content', '')}\n\n用户问题: {question}"
                
                # 使用MiniRAG查询
                response = await rag.aquery(image_context, param=QueryParam(mode=mode))
                
                return {
                    "status": "success",
                    "question": question,
                    "answer": response,
                    "image_content": parsed_result.get('content', ''),
                    "image_info": {
                        "size": parsed_result.get('image_size', ''),
                        "mode": parsed_result.get('image_mode', '')
                    }
                }
            finally:
                # 清理临时文件
                if temp_image_path.exists():
                    temp_image_path.unlink()
                    
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"图片问答失败: {str(e)}")

    @app.get("/files/status", dependencies=[Depends(optional_api_key)])
    async def get_files_status():
        """获取文件处理状态"""
        try:
            # 扫描工作目录中的文件
            working_dir = Path(args.working_dir)
            files_info = []
            
            if working_dir.exists():
                for file_path in working_dir.rglob("*"):
                    if file_path.is_file():
                        try:
                            stat = file_path.stat()
                            files_info.append({
                                "name": file_path.name,
                                "path": str(file_path.relative_to(working_dir)),
                                "size": stat.st_size,
                                "modified": stat.st_mtime,
                                "type": file_path.suffix.lower()
                            })
                        except Exception:
                            continue
            
            return {
                "status": "success",
                "total_files": len(files_info),
                "files": files_info
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"获取文件状态失败: {str(e)}")

    @app.get("/supported-formats", dependencies=[Depends(optional_api_key)])
    async def get_supported_formats():
        """获取支持的文件格式"""
        try:
            from minirag.file_parser import FileParserManager
            parser_manager = FileParserManager()
            formats = parser_manager.get_supported_formats()
            
            return {
                "status": "success",
                "supported_formats": formats
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"获取支持格式失败: {str(e)}")

    @app.delete("/files/{file_name}", dependencies=[Depends(optional_api_key)])
    async def delete_file(file_name: str):
        """删除指定文件"""
        try:
            working_dir = Path(args.working_dir)
            file_path = working_dir / file_name
            
            if not file_path.exists():
                raise HTTPException(status_code=404, detail="文件不存在")
            
            # 从MiniRAG中移除相关内容（如果已索引）
            # 这里可以添加从向量数据库、知识图谱等存储中移除内容的逻辑
            
            # 删除文件
            file_path.unlink()
            
            return {
                "status": "success",
                "message": f"文件 {file_name} 删除成功"
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"删除文件失败: {str(e)}")

    @app.get("/system/info", dependencies=[Depends(optional_api_key)])
    async def get_system_info():
        """获取系统详细信息"""
        try:
            import psutil
            import platform
            
            # 系统信息
            system_info = {
                "platform": platform.system(),
                "platform_version": platform.version(),
                "python_version": platform.python_version(),
                "architecture": platform.architecture()[0]
            }
            
            # 内存信息
            memory = psutil.virtual_memory()
            memory_info = {
                "total": memory.total,
                "available": memory.available,
                "percent": memory.percent,
                "used": memory.used
            }
            
            # 磁盘信息
            disk = psutil.disk_usage(str(args.working_dir))
            disk_info = {
                "total": disk.total,
                "used": disk.used,
                "free": disk.free,
                "percent": disk.percent
            }
            
            return {
                "status": "success",
                "system": system_info,
                "memory": memory_info,
                "disk": disk_info,
                "working_directory": str(args.working_dir),
                "input_directory": str(args.input_dir)
            }
        except ImportError:
            # 如果没有psutil，返回基本信息
            return {
                "status": "success",
                "working_directory": str(args.working_dir),
                "input_directory": str(args.input_dir),
                "note": "安装psutil可获取更详细的系统信息"
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"获取系统信息失败: {str(e)}")

    @app.post("/files/batch-process", dependencies=[Depends(optional_api_key)])
    async def batch_process_files(
        background_tasks: BackgroundTasks,
        directory: str = Form(...),
        file_types: str = Form("all")
    ):
        """批量处理指定目录中的文件"""
        try:
            target_dir = Path(directory)
            if not target_dir.exists() or not target_dir.is_dir():
                raise HTTPException(status_code=400, detail="目录不存在或不是有效目录")
            
            # 获取文件列表
            file_extensions = []
            if file_types != "all":
                file_extensions = [ext.strip() for ext in file_types.split(",")]
            
            files_to_process = []
            for file_path in target_dir.rglob("*"):
                if file_path.is_file():
                    if not file_extensions or file_path.suffix.lower() in file_extensions:
                        files_to_process.append(str(file_path))
            
            if not files_to_process:
                return {
                    "status": "success",
                    "message": "没有找到符合条件的文件",
                    "total_files": 0
                }
            
            # 后台任务：批量处理文件
            async def process_batch_files():
                try:
                    from minirag.file_parser import parse_multiple_files
                    
                    # 批量解析文件
                    parsed_results = await parse_multiple_files(files_to_process)
                    
                    # 统计结果
                    success_count = 0
                    error_count = 0
                    
                    for result in parsed_results:
                        if 'error' not in result:
                            try:
                                # 插入到MiniRAG
                                await rag.ainsert(result.get('content', ''))
                                success_count += 1
                            except Exception as e:
                                logging.error(f"插入文件内容失败: {result.get('file_path', 'unknown')}: {e}")
                                error_count += 1
                        else:
                            error_count += 1
                    
                    logging.info(f"批量处理完成: 成功 {success_count} 个文件, 失败 {error_count} 个文件")
                    
                except Exception as e:
                    logging.error(f"批量处理文件失败: {e}")
            
            background_tasks.add_task(process_batch_files)
            
            return {
                "status": "success",
                "message": f"开始批量处理 {len(files_to_process)} 个文件",
                "total_files": len(files_to_process),
                "file_types": file_extensions if file_extensions else "all"
            }
            
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"批量处理失败: {str(e)}")

    @app.get("/files/search", dependencies=[Depends(optional_api_key)])
    async def search_files(
        query: str,
        file_type: str = None,
        size_min: int = None,
        size_max: int = None
    ):
        """搜索文件"""
        try:
            working_dir = Path(args.working_dir)
            files_info = []
            
            if working_dir.exists():
                for file_path in working_dir.rglob("*"):
                    if file_path.is_file():
                        try:
                            stat = file_path.stat()
                            
                            # 检查文件类型过滤
                            if file_type and file_path.suffix.lower() != file_type.lower():
                                continue
                            
                            # 检查文件大小过滤
                            if size_min and stat.st_size < size_min:
                                continue
                            if size_max and stat.st_size > size_max:
                                continue
                            
                            # 检查文件名是否匹配查询
                            if query.lower() not in file_path.name.lower():
                                continue
                            
                            files_info.append({
                                "name": file_path.name,
                                "path": str(file_path.relative_to(working_dir)),
                                "size": stat.st_size,
                                "modified": stat.st_mtime,
                                "type": file_path.suffix.lower()
                            })
                        except Exception:
                            continue
            
            return {
                "status": "success",
                "query": query,
                "total_files": len(files_info),
                "files": files_info
            }
            
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"搜索文件失败: {str(e)}")

    @app.get("/api/endpoints", dependencies=[Depends(optional_api_key)])
    async def get_api_endpoints():
        """获取所有可用的API端点信息"""
        endpoints = [
            {
                "path": "/v1/chat/completions",
                "method": "POST",
                "description": "智能问答接口",
                "parameters": ["model", "messages", "mode", "stream"]
            },
            {
                "path": "/upload/file",
                "method": "POST",
                "description": "上传并解析文档文件",
                "parameters": ["file", "description"]
            },
            {
                "path": "/upload/image",
                "method": "POST",
                "description": "上传并解析图片（多模态理解）",
                "parameters": ["image", "description"]
            },
            {
                "path": "/query/image",
                "method": "POST",
                "description": "基于图片的问答",
                "parameters": ["image", "question", "mode"]
            },
            {
                "path": "/files/status",
                "method": "GET",
                "description": "获取文件处理状态",
                "parameters": []
            },
            {
                "path": "/supported-formats",
                "method": "GET",
                "description": "获取支持的文件格式",
                "parameters": []
            },
            {
                "path": "/documents",
                "method": "GET",
                "description": "获取已索引的文档列表",
                "parameters": []
            },
            {
                "path": "/health",
                "method": "GET",
                "description": "获取系统状态",
                "parameters": []
            },
            {
                "path": "/system/info",
                "method": "GET",
                "description": "获取系统详细信息",
                "parameters": []
            },
            {
                "path": "/files/batch-process",
                "method": "POST",
                "description": "批量处理文件",
                "parameters": ["directory", "file_types"]
            },
            {
                "path": "/files/search",
                "method": "GET",
                "description": "搜索文件",
                "parameters": ["query", "file_type", "size_min", "size_max"]
            }
        ]
        
        return {
            "status": "success",
            "total_endpoints": len(endpoints),
            "endpoints": endpoints
        }

    # webui mount /webui/index.html
    # app.mount(
    #     "/webui",
    #     StaticFiles(
    #         directory=Path(__file__).resolve().parent / "webui" / "static", html=True
    #     ),
    #     name="webui_static",
    # )

    # Serve the static files
    static_dir = Path(__file__).parent / "static"
    static_dir.mkdir(exist_ok=True)
    app.mount("/", StaticFiles(directory=static_dir, html=True), name="static")

    return app


def main():
    args = parse_args()
    import uvicorn

    app = create_app(args)
    display_splash_screen(args)
    uvicorn_config = {
        "app": app,
        "host": args.host,
        "port": args.port,
    }
    if args.ssl:
        uvicorn_config.update(
            {
                "ssl_certfile": args.ssl_certfile,
                "ssl_keyfile": args.ssl_keyfile,
            }
        )
    uvicorn.run(**uvicorn_config)


if __name__ == "__main__":
    main()
