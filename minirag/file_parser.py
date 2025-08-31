"""
文件解析器模块 - 支持多种文件格式的解析
支持格式：PDF、PPT、Word、Markdown、TXT等
"""

import os
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
from abc import ABC, abstractmethod
import asyncio

# 文件解析相关依赖
try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

try:
    from docx import Document
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

logger = logging.getLogger(__name__)


class FileParser(ABC):
    """文件解析器基类"""
    
    @abstractmethod
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析文件并返回结构化数据"""
        pass
    
    @abstractmethod
    def can_parse(self, file_path: str) -> bool:
        """检查是否可以解析该文件"""
        pass


class TextFileParser(FileParser):
    """文本文件解析器"""
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为文本文件"""
        text_extensions = {'.txt', '.md', '.markdown', '.log', '.csv', '.json', '.xml', '.html', '.htm'}
        return Path(file_path).suffix.lower() in text_extensions
    
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'content': content,
                'file_type': 'text',
                'file_path': file_path,
                'file_size': len(content),
                'encoding': 'utf-8'
            }
        except Exception as e:
            logger.error(f"解析文本文件失败 {file_path}: {e}")
            return {'error': str(e), 'file_path': file_path}


class PDFFileParser(FileParser):
    """PDF文件解析器"""
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为PDF文件"""
        return Path(file_path).suffix.lower() == '.pdf' and PDF_AVAILABLE
    
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析PDF文件"""
        try:
            logger.info(f"开始解析PDF文件: {file_path}")
            logger.info(f"文件路径类型: {type(file_path)}")
            logger.info(f"文件是否存在: {Path(file_path).exists()}")
            logger.info(f"文件大小: {Path(file_path).stat().st_size if Path(file_path).exists() else 'N/A'}")
            
            content = ""
            
            # 尝试使用pdfplumber（更好的文本提取）
            if pdfplumber:
                logger.info("尝试使用pdfplumber解析PDF")
                try:
                    with pdfplumber.open(file_path) as pdf:
                        logger.info(f"PDF页数: {len(pdf.pages)}")
                        for i, page in enumerate(pdf.pages):
                            logger.info(f"正在处理第 {i+1} 页")
                            page_text = page.extract_text()
                            if page_text:
                                content += page_text + "\n"
                                logger.info(f"第 {i+1} 页提取文本长度: {len(page_text)} 字符")
                            else:
                                logger.warning(f"第 {i+1} 页没有提取到文本")
                except Exception as e:
                    logger.error(f"pdfplumber解析失败: {e}")
            
            # 如果pdfplumber失败，尝试PyPDF2
            if not content and PyPDF2:
                logger.info("pdfplumber失败，尝试使用PyPDF2解析PDF")
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        logger.info(f"PyPDF2读取成功，页数: {len(pdf_reader.pages)}")
                        for i, page in enumerate(pdf_reader.pages):
                            logger.info(f"正在处理第 {i+1} 页")
                            page_text = page.extract_text()
                            if page_text:
                                content += page_text + "\n"
                                logger.info(f"第 {i+1} 页提取文本长度: {len(page_text)} 字符")
                            else:
                                logger.warning(f"第 {i+1} 页没有提取到文本")
                except Exception as e:
                    logger.error(f"PyPDF2解析失败: {e}")
            
            if not content:
                logger.error("所有PDF解析方法都失败了，无法提取文本内容")
                raise Exception("无法提取PDF文本内容")
            
            logger.info(f"PDF解析成功，总内容长度: {len(content)} 字符")
            return {
                'content': content.strip(),
                'file_type': 'pdf',
                'file_path': file_path,
                'file_size': len(content),
                'page_count': len(content.split('\n'))
            }
        except Exception as e:
            logger.error(f"解析PDF文件失败 {file_path}: {e}")
            return {'error': str(e), 'file_path': file_path}


class PowerPointParser(FileParser):
    """PowerPoint文件解析器"""
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为PPT文件"""
        ppt_extensions = {'.pptx', '.ppt'}
        return Path(file_path).suffix.lower() in ppt_extensions and PPT_AVAILABLE
    
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析PowerPoint文件"""
        try:
            prs = Presentation(file_path)
            content = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content += f"\n--- 幻灯片 {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text.strip() + "\n"
            
            return {
                'content': content.strip(),
                'file_type': 'powerpoint',
                'file_path': file_path,
                'file_size': len(content),
                'slide_count': len(prs.slides)
            }
        except Exception as e:
            logger.error(f"解析PowerPoint文件失败 {file_path}: {e}")
            return {'error': str(e), 'file_path': file_path}


class WordDocumentParser(FileParser):
    """Word文档解析器"""
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为Word文档"""
        word_extensions = {'.docx', '.doc'}
        return Path(file_path).suffix.lower() in word_extensions and WORD_AVAILABLE
    
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析Word文档"""
        try:
            doc = Document(file_path)
            content = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text.strip() + "\n"
            
            # 提取表格内容
            for table in doc.tables:
                content += "\n--- 表格 ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    content += row_text + "\n"
            
            return {
                'content': content.strip(),
                'file_type': 'word',
                'file_path': file_path,
                'file_size': len(content),
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()])
            }
        except Exception as e:
            logger.error(f"解析Word文档失败 {file_path}: {e}")
            return {'error': str(e), 'file_path': file_path}


class ImageParser(FileParser):
    """图片文件解析器（多模态理解）"""
    
    def __init__(self, multimodal_model_func=None):
        """初始化图片解析器
        
        Args:
            multimodal_model_func: 多模态模型函数，用于理解图片内容
        """
        self.multimodal_model_func = multimodal_model_func
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为图片文件"""
        image_extensions = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif', '.webp'}
        return Path(file_path).suffix.lower() in image_extensions and IMAGE_AVAILABLE
    
    async def parse(self, file_path: str, description: str = "") -> Dict[str, Any]:
        """解析图片文件（多模态理解）"""
        try:
            # 打开图片
            image = Image.open(file_path)
            
            # 如果有多模态模型，使用模型理解图片
            if self.multimodal_model_func:
                try:
                    # 使用多模态模型理解图片
                    image_understanding = await self._understand_image_with_model(image, description)
                    content = image_understanding
                except Exception as model_error:
                    logger.warning(f"多模态模型理解失败，回退到OCR: {model_error}")
                    # 回退到OCR
                    content = await self._extract_text_with_ocr(image)
            else:
                # 没有多模态模型时，使用OCR
                content = await self._extract_text_with_ocr(image)
            
            return {
                'content': content.strip(),
                'file_type': 'image',
                'file_path': file_path,
                'file_size': len(content),
                'image_size': image.size,
                'image_mode': image.mode,
                'processing_method': 'multimodal' if self.multimodal_model_func else 'ocr'
            }
        except Exception as e:
            logger.error(f"解析图片文件失败 {file_path}: {e}")
            return {'error': str(e), 'file_path': file_path}
    
    async def _understand_image_with_model(self, image: Image.Image, description: str = "") -> str:
        """使用多模态模型理解图片"""
        try:
            # 构建提示词
            if description:
                prompt = f"请详细描述这张图片的内容。用户描述：{description}\n\n请从以下方面进行分析：\n1. 图片的主要内容和场景\n2. 可见的文字、图表或数据\n3. 图片的整体含义和用途\n4. 任何重要的细节或特征"
            else:
                prompt = "请详细描述这张图片的内容，包括主要场景、可见的文字、图表、数据以及图片的整体含义。"
            
            # 调用多模态模型
            if callable(self.multimodal_model_func):
                # 如果是函数，直接调用
                if asyncio.iscoroutinefunction(self.multimodal_model_func):
                    response = await self.multimodal_model_func(image, prompt)
                else:
                    response = self.multimodal_model_func(image, prompt)
            else:
                # 如果是对象，调用其方法
                if hasattr(self.multimodal_model_func, 'aquery'):
                    response = await self.multimodal_model_func.aquery(prompt, image=image)
                elif hasattr(self.multimodal_model_func, 'query'):
                    response = self.multimodal_model_func.query(prompt, image=image)
                else:
                    raise ValueError("多模态模型对象必须具有 query 或 aquery 方法")
            
            return str(response)
            
        except Exception as e:
            logger.error(f"多模态模型理解图片失败: {e}")
            raise
    
    async def _extract_text_with_ocr(self, image: Image.Image) -> str:
        """使用OCR提取图片中的文字"""
        try:
            # 使用OCR提取文本
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            return text
        except Exception as e:
            logger.error(f"OCR提取文字失败: {e}")
            return f"图片OCR识别失败: {str(e)}"


class FileParserManager:
    """文件解析器管理器"""
    
    def __init__(self, multimodal_model_func=None):
        """初始化解析器管理器
        
        Args:
            multimodal_model_func: 多模态模型函数或对象，用于理解图片内容
        """
        self.multimodal_model_func = multimodal_model_func
        self.parsers: List[FileParser] = [
            TextFileParser(),
            PDFFileParser(),
            PowerPointParser(),
            WordDocumentParser(),
            ImageParser(multimodal_model_func)
        ]
        
        # 记录支持的格式
        self.supported_formats = self._get_supported_formats()
    
    def _get_supported_formats(self) -> Dict[str, List[str]]:
        """获取支持的格式列表"""
        formats = {}
        for parser in self.parsers:
            if isinstance(parser, TextFileParser):
                formats['text'] = ['.txt', '.md', '.markdown', '.log', '.csv', '.json', '.xml', '.html', '.htm']
            elif isinstance(parser, PDFFileParser) and PDF_AVAILABLE:
                formats['pdf'] = ['.pdf']
            elif isinstance(parser, PowerPointParser) and PPT_AVAILABLE:
                formats['powerpoint'] = ['.pptx', '.ppt']
            elif isinstance(parser, WordDocumentParser) and WORD_AVAILABLE:
                formats['word'] = ['.docx', '.doc']
            elif isinstance(parser, ImageParser) and IMAGE_AVAILABLE:
                formats['image'] = ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']
        
        return formats
    
    def get_parser(self, file_path: str) -> Optional[FileParser]:
        """根据文件路径获取合适的解析器"""
        for parser in self.parsers:
            if parser.can_parse(file_path):
                return parser
        return None
    
    async def parse_file(self, file_path: str) -> Dict[str, Any]:
        """解析单个文件"""
        parser = self.get_parser(file_path)
        if not parser:
            return {
                'error': f'不支持的文件格式: {Path(file_path).suffix}',
                'file_path': file_path
            }
        
        return await parser.parse(file_path)
    
    async def parse_files(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """批量解析文件"""
        tasks = [self.parse_file(file_path) for file_path in file_paths]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # 处理异常结果
        parsed_results = []
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                parsed_results.append({
                    'error': str(result),
                    'file_path': file_paths[i]
                })
            else:
                parsed_results.append(result)
        
        return parsed_results
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """获取支持的格式信息"""
        return self.supported_formats.copy()


# 便捷函数
async def parse_single_file(file_path: str, multimodal_model_func=None, description: str = "") -> Dict[str, Any]:
    """解析单个文件的便捷函数
    
    Args:
        file_path: 文件路径
        multimodal_model_func: 多模态模型函数或对象
        description: 图片描述（仅对图片文件有效）
    """
    import time
    start_time = time.time()
    
    logger.info(f"开始解析文件: {file_path}")
    
    # 创建解析器管理器
    manager_start = time.time()
    manager = FileParserManager(multimodal_model_func)
    manager_time = time.time() - manager_start
    logger.info(f"解析器管理器创建完成 (耗时: {manager_time:.3f}s)")
    
    # 解析文件
    parse_start = time.time()
    result = await manager.parse_file(file_path)
    parse_time = time.time() - parse_start
    
    if 'error' in result:
        total_time = time.time() - start_time
        logger.error(f"文件解析失败: {file_path} (解析耗时: {parse_time:.3f}s, 总耗时: {total_time:.3f}s)")
        logger.error(f"   错误信息: {result['error']}")
        return result
    
    # 如果是图片文件且有描述，重新解析以包含描述
    if result.get('file_type') == 'image' and description and multimodal_model_func:
        logger.info(f"图片文件，开始重新解析以包含描述")
        image_parse_start = time.time()
        image_parser = manager.get_parser(file_path)
        if isinstance(image_parser, ImageParser):
            result = await image_parser.parse(file_path, description)
        image_parse_time = time.time() - image_parse_start
        logger.info(f"图片重新解析完成 (耗时: {image_parse_time:.3f}s)")
    
    total_time = time.time() - start_time
    content_length = len(result.get('content', ''))
    logger.info(f"文件解析完成: {file_path}")
    logger.info(f"   - 文件类型: {result.get('file_type', 'unknown')}")
    logger.info(f"   - 内容长度: {content_length} 字符")
    logger.info(f"   - 解析耗时: {parse_time:.3f}s")
    logger.info(f"   - 总耗时: {total_time:.3f}s")
    
    return result


async def parse_multiple_files(file_paths: List[str], multimodal_model_func=None) -> List[Dict[str, Any]]:
    """解析多个文件的便捷函数
    
    Args:
        file_paths: 文件路径列表
        multimodal_model_func: 多模态模型函数或对象
    """
    manager = FileParserManager(multimodal_model_func)
    return await manager.parse_files(file_paths)
